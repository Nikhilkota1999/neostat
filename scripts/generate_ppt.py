#!/usr/bin/env python3
"""
Generate the case study PPT deck for the NeoStats AI Engineer challenge.

Creates: deliverables/NeoStats_Case_Study.pptx
"""
import os
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except Exception as e:
    raise SystemExit(
        "python-pptx is required. Install with: python3 -m pip install python-pptx\n"
        f"Import error: {e}"
    )


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullets_slide(prs, heading, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = bullet
        p.level = 0
    return slide


def main():
    out_dir = os.path.join("deliverables")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "NeoStats_Case_Study.pptx")

    today = datetime.now().strftime("%Y-%m-%d")
    prs = Presentation()

    # Title
    add_title_slide(
        prs,
        title="NeoStats AI Engineer Case Study",
        subtitle=f"The Chatbot Blueprint — {today}"
    )

    # Use Case Objective
    add_bullets_slide(
        prs,
        "Use Case Objective",
        [
            "Build an intelligent chatbot with contextual answers",
            "Support local knowledge via RAG and live web search",
            "Deliver concise or detailed responses on demand",
        ],
    )

    # Approach
    add_bullets_slide(
        prs,
        "Approach",
        [
            "Start from Streamlit template and modularize by domain",
            "Implement RAG: load → split → embed → retrieve",
            "Add web search tool for real-time context",
            "Make response mode a first-class UI control",
            "Centralize configuration and keys in config",
        ],
    )

    # Solution Architecture
    add_bullets_slide(
        prs,
        "Solution Architecture",
        [
            "UI: Streamlit chat (app.py)",
            "LLMs: Mistral (default) + Groq/Gemini/OpenRouter supported",
            "RAG: Chroma vector store + HuggingFace (or OpenAI) embeddings",
            "Search: DuckDuckGo (ddgs) for lightweight live web results",
            "Config: .env → config/config.py (no keys in code)",
        ],
    )

    # Features Implemented
    add_bullets_slide(
        prs,
        "Features Implemented",
        [
            "Concise vs Detailed response modes",
            "Local document upload and indexing (PDF/Text)",
            "Top-k retrieval with source snippets",
            "Optional live web search context",
            "LLM health check and error handling",
        ],
    )

    # Technical Details
    add_bullets_slide(
        prs,
        "Technical Details",
        [
            "Vector store: Chroma (in-memory by default)",
            "Embeddings: sentence-transformers/all-MiniLM-L6-v2 (default)",
            "LLM provider set via LLM_PROVIDER in .env",
            "Extensible utils: utils/rag.py and utils/search.py",
        ],
    )

    # Challenges
    add_bullets_slide(
        prs,
        "Challenges",
        [
            "Provider compatibility and versioning across SDKs",
            "Graceful fallbacks when keys/models are unavailable",
            "Efficient chunking and retrieval quality for mixed PDFs",
        ],
    )

    # Deployment
    add_bullets_slide(
        prs,
        "Deployment",
        [
            "Deploy to Streamlit Cloud",
            "Set API keys as secrets (no repo commits)",
            "Command: streamlit run AI_UseCase/app.py",
            "App link: <add Streamlit Cloud URL here>",
        ],
    )

    # Next Steps
    add_bullets_slide(
        prs,
        "Next Steps",
        [
            "Persist Chroma index to disk (persist_directory)",
            "Add sidebar model selector for quick switching",
            "Inline citations matching sources panel",
        ],
    )

    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()

